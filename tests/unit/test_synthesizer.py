"""Offline unit tests for ria/synthesizer.py's local file generation (no API calls).

DOCX/PPTX generation is local file I/O, not network -- genuinely offline-testable. Tests
write into pytest's tmp_path, not the real outputs/ directory (an absolute path joined onto
PROJECT_ROOT via pathlib resolves to the absolute path itself, discarding PROJECT_ROOT).
"""

from datetime import date

from docx import Document
from pptx import Presentation

from ria import synthesizer
from ria.models import RegulatoryDocument


def _doc(**overrides):
    base = dict(
        document_number="2026-00001",
        title="Test Rule Requiring Executive Review",
        document_type="Rule",
        agencies=["Food and Drug Administration"],
        publication_date=date(2026, 7, 1),
        html_url="https://example.gov/doc",
    )
    base.update(overrides)
    return RegulatoryDocument(**base)


def _settings(tmp_path):
    class _Settings:
        output = {
            "output_path_docx": str(tmp_path / "docx"),
            "output_path_pptx": str(tmp_path / "pptx"),
            "docx_template": "",
            "pptx_template": "",
        }
    return _Settings()


_BRIEFING = {
    "executive_summary": "This rule requires updating our drug registration process.",
    "remediation_plan": [
        {"action": "Update registration SOP", "owner": "Regulatory Affairs",
         "due_date": "2026-08-01", "priority": "high"},
        {"action": "Train quality team", "owner": "Quality Assurance",
         "due_date": "2026-09-01", "priority": "medium"},
    ],
}
_MATERIALITY = {"impact_score": 65, "risk_level": "high"}
_EVALUATOR_DECISION = {
    "autonomy_tier": 2, "execute": False, "escalate": False,
    "scores": {"overall_confidence": 0.8},
}


def test_resolve_template_returns_none_when_unconfigured(tmp_path):
    settings = _settings(tmp_path)
    assert synthesizer._resolve_template(settings, "docx_template") is None


def test_resolve_template_returns_none_when_configured_path_missing(tmp_path):
    class _Settings:
        output = {"docx_template": "config/does_not_exist.docx"}
    assert synthesizer._resolve_template(_Settings(), "docx_template") is None


def test_output_path_creates_directory_and_names_by_document_number(tmp_path):
    settings = _settings(tmp_path)
    path = synthesizer._output_path(settings, "docx", _doc())
    assert path.parent.exists()
    assert path.name == "2026-00001.docx"


def test_write_docx_produces_a_valid_readable_document(tmp_path):
    settings = _settings(tmp_path)
    path = synthesizer._write_docx(_doc(), _BRIEFING, _MATERIALITY, _EVALUATOR_DECISION, settings)
    assert path.exists()

    reopened = Document(str(path))
    full_text = "\n".join(p.text for p in reopened.paragraphs)
    assert "Test Rule Requiring Executive Review" in full_text
    assert "This rule requires updating our drug registration process." in full_text
    assert "Tier 2" in full_text


def test_write_docx_table_has_one_row_per_remediation_item(tmp_path):
    settings = _settings(tmp_path)
    path = synthesizer._write_docx(_doc(), _BRIEFING, _MATERIALITY, _EVALUATOR_DECISION, settings)
    reopened = Document(str(path))
    table = reopened.tables[0]
    assert len(table.rows) == 1 + len(_BRIEFING["remediation_plan"])  # header + 2 items
    assert table.rows[1].cells[0].text == "Update registration SOP"


def test_write_pptx_produces_a_valid_readable_presentation(tmp_path):
    settings = _settings(tmp_path)
    path = synthesizer._write_pptx(_doc(), _BRIEFING, _MATERIALITY, _EVALUATOR_DECISION, settings)
    assert path.exists()

    reopened = Presentation(str(path))
    assert len(reopened.slides) == 4  # title + summary + risk + plan

    titles = [s.shapes.title.text for s in reopened.slides if s.shapes.title is not None]
    assert "Remediation Plan" in titles


def test_write_pptx_handles_empty_remediation_plan(tmp_path):
    settings = _settings(tmp_path)
    empty_briefing = {**_BRIEFING, "remediation_plan": []}
    path = synthesizer._write_pptx(_doc(), empty_briefing, _MATERIALITY, _EVALUATOR_DECISION, settings)
    assert path.exists()  # should not crash on an empty plan


def test_scrub_removes_cfr_and_usc_citations():
    dirty = "This is required under 21 CFR Part 207 and 21 U.S.C. 333(f)."
    cleaned = synthesizer._scrub_executive_summary(dirty)
    assert "CFR" not in cleaned
    assert "U.S.C." not in cleaned
    assert "applicable federal regulations" in cleaned


def test_scrub_replaces_jargon_terms():
    dirty = "The FDA considers our product misbranded, notwithstanding prior corrective action."
    cleaned = synthesizer._scrub_executive_summary(dirty)
    assert "misbranded" not in cleaned.lower()
    assert "notwithstanding" not in cleaned.lower()
    assert "mislabeled" in cleaned.lower()


def test_scrub_leaves_clean_text_unchanged():
    clean = "We need to fix our sterility controls within two weeks."
    assert synthesizer._scrub_executive_summary(clean) == clean


def test_scrub_logs_only_when_something_actually_changed():
    calls = []

    class _Logger:
        def info(self, msg):
            calls.append(msg)

    synthesizer._scrub_executive_summary("Nothing to fix here.", _Logger(), "DOC-1")
    assert calls == []

    synthesizer._scrub_executive_summary("Governed by 21 CFR Part 207.", _Logger(), "DOC-1")
    assert len(calls) == 1
    assert "jargon_scrub" in calls[0]
